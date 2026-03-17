"""Advanced image understanding using multimodal LLMs with SiliconFlow support."""
 
import base64
import io
import json
import logging
import requests
from pathlib import Path
from typing import Any, Union

import tempfile
import uuid
import zipfile
import shutil

import contextlib
from typing import List
 
from langchain_core.messages import HumanMessage

from logging.handlers import RotatingFileHandler
import asyncio
 
# Type alias for image input
ImagePathOrImage = Union[Path, 'PIL.Image.Image']

DEFAULT_LANGUAGE = "zh"  # 默认语言：zh(中文) 或 en(英文)
def get_language() -> str:
    """
    获取当前配置的语言
    可以通过环境变量或配置文件修改
    """
    import os
    return os.getenv("IMAGE_UNDERSTANDING_LANG", DEFAULT_LANGUAGE)
 
def setup_image_logger():
    log_dir = Path(__file__).parent.parent.parent / "logs"
    log_dir.mkdir(exist_ok=True)
    log_file = log_dir / "image_understanding.log"
    
    handler = RotatingFileHandler(
        log_file,
        maxBytes=10*1024*1024,
        backupCount=5,
        encoding='utf-8'
    )
    handler.setFormatter(logging.Formatter(
        "%(asctime)s - %(name)s - %(levelname)s - %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S"
    ))
    
    logger = logging.getLogger(__name__)
    logger.setLevel(logging.INFO)
    logger.addHandler(handler)
    logger.propagate = True  # 同时也输出到其他地方
    
    return logger

logger = setup_image_logger()

@contextlib.asynccontextmanager
async def temp_image_manager(image_paths: List[Path]):
    """
    上下文管理器：确保临时图片文件被清理
    """
    try:
        yield image_paths
    finally:
        # 清理所有临时文件
        for img_path in image_paths:
            if isinstance(img_path, Path) and img_path.exists():
                try:
                    # 检查是否是临时目录中的文件
                    if str(img_path).startswith(tempfile.gettempdir()):
                        img_path.unlink(missing_ok=True)
                        logger.info(f"Cleaned up temp file: {img_path.name}")
                except Exception as e:
                    logger.warning(f"Failed to cleanup {img_path}: {e}")
        
        # 清理空目录
        if image_paths:
            try:
                parent = image_paths[0].parent
                if str(parent).startswith(tempfile.gettempdir()) and parent.exists():
                    # 检查目录是否为空
                    if not any(parent.iterdir()):
                        shutil.rmtree(parent, ignore_errors=True)
                        logger.info(f"Cleaned up temp dir: {parent}")
            except Exception as e:
                logger.warning(f"Failed to cleanup temp dir: {e}")

async def extract_ppt_images_safe(pptx_path: Path) -> list[Path]:
    """
    安全提取PPT图片，处理python-pptx的各种bug
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    temp_dir = Path(tempfile.mkdtemp())
    extracted_images = []
    failed_extractions = []
    processed_hashes = set()  # 用于去重

    def get_image_hash(image_bytes: bytes) -> str:
        """计算图片内容的哈希值用于去重"""
        import hashlib
        return hashlib.md5(image_bytes).hexdigest()
    
    def add_image_if_unique(image_bytes: bytes, img_path: Path) -> bool:
        """如果图片未重复则添加"""
        img_hash = get_image_hash(image_bytes)
        if img_hash in processed_hashes:
            logger.info(f"跳过重复图片: {img_path.name}")
            return False
        processed_hashes.add(img_hash)
        img_path.write_bytes(image_bytes)
        extracted_images.append(img_path)
        return True
    
    try:
        prs = Presentation(pptx_path)
        total_shapes = 0
        picture_shapes = 0
        actual_images = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                total_shapes += 1
                
                # 更严格的图片检测
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_shapes += 1
                    
                    # 检查是否有有效的 blip 元素
                    try:
                        if not hasattr(shape, '_pic') or shape._pic is None:
                            logger.debug(f"Slide {slide_num}: 形状无 _pic 属性，跳过")
                            continue
                        
                        blip = shape._pic.spPr.blipFill.blip
                        if blip is None:
                            logger.debug(f"Slide {slide_num}: blip 为空，跳过")
                            continue
                        
                        # 检查是否是嵌入图片（不是链接）
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if not embed:
                            logger.debug(f"Slide {slide_num}: 链接图片，跳过")
                            continue
                        
                    except Exception as e:
                        logger.debug(f"Slide {slide_num}: 图片检测失败: {e}")
                        continue
                    
                    # 尝试提取
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        if not image_bytes or len(image_bytes) < 100:  # 过滤太小的文件
                            logger.warning(f"Slide {slide_num}: 图片内容太小或为空")
                            continue
                        
                        ext = image.ext if image.ext else 'png'
                        img_path = temp_dir / f"slide{slide_num}_{uuid.uuid4().hex}.{ext}"
                        
                        if add_image_if_unique(image_bytes, img_path):
                            actual_images += 1
                            logger.info(f"成功提取图片: {img_path.name} ({len(image_bytes)} bytes)")
                        
                    except AttributeError as e:
                        if "'Part' object has no attribute 'image'" in str(e):
                            logger.warning(f"Slide {slide_num}: JPEG/python-pptx bug")
                            failed_extractions.append({'slide': slide_num, 'reason': 'jpeg_bug'})
                        else:
                            logger.warning(f"Slide {slide_num}: 属性错误 - {e}")
                    except Exception as e:
                        logger.warning(f"Slide {slide_num}: 提取失败 - {type(e).__name__}: {e}")
        
        logger.info(f"PPT分析完成: 总形状={total_shapes}, 图片形状={picture_shapes}, "
                   f"成功提取={len(extracted_images)}, 失败={len(failed_extractions)}")
        
        # 如果python-pptx没有提取到所有图片，尝试zipfile方法
        if len(extracted_images) < picture_shapes:
            logger.info(f"尝试使用zipfile直接提取图片... "
                       f"(已提取{len(extracted_images)}/{picture_shapes})")
            zip_images = await extract_images_via_zip(pptx_path, temp_dir)
            # 去重：避免重复提取已成功的图片
            existing_names = {p.name for p in extracted_images}
            for img_path in zip_images:
                if img_path.name not in existing_names:
                    extracted_images.append(img_path)
                    logger.info(f"zipfile补充提取: {img_path.name}")
        
    except Exception as e:
        logger.error(f"解析PPT失败: {e}")
        extracted_images = await extract_images_via_zip(pptx_path, temp_dir)
    
    logger.info(f"最终提取图片数量: {len(extracted_images)}")
    return extracted_images

async def extract_docx_images_safe(docx_path: Path) -> list[Path]:
    """
    安全提取DOCX图片，处理各种格式问题
    """
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    
    temp_dir = Path(tempfile.mkdtemp())
    extracted_images = []
    
    try:
        doc = Document(docx_path)
        
        # 方法1：通过关系提取
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    if not hasattr(image_part, 'blob'):
                        logger.warning(f"关系 {rel.target_ref} 不是图片 part，跳过")
                        continue
                    
                    image_bytes = image_part.blob
                    if not image_bytes:
                        logger.warning(f"图片 {rel.target_ref} 内容为空，跳过")
                        continue
                    
                    # 尝试确定图片格式
                    try:
                        from PIL import Image
                        img = Image.open(io.BytesIO(image_bytes))
                        ext = img.format.lower() if img.format else 'png'
                    except:
                        ext = 'png'
                    
                    img_path = temp_dir / f"docx_image_{uuid.uuid4().hex}.{ext}"
                    img_path.write_bytes(image_bytes)
                    extracted_images.append(img_path)
                    logger.info(f"成功提取DOCX图片: {img_path.name}")
                    
                except AttributeError as e:
                    logger.warning(f"提取图片 {rel.target_ref} 属性错误: {e}")
                except Exception as e:
                    logger.warning(f"提取图片 {rel.target_ref} 失败: {type(e).__name__}: {e}")
        
        # 方法2：如果方法1没有提取到图片，尝试zipfile
        if not extracted_images:
            logger.info("尝试使用zipfile直接提取DOCX图片...")
            extracted_images = await extract_docx_images_via_zip(docx_path, temp_dir)
            
    except Exception as e:
        logger.error(f"解析DOCX失败: {e}")
        # 尝试zipfile作为备用
        extracted_images = await extract_docx_images_via_zip(docx_path, temp_dir)
    
    return extracted_images


async def extract_docx_images_via_zip(docx_path: Path, output_dir: Path) -> list[Path]:
    """
    使用zipfile直接解压DOCX中的图片（备用方法）
    """
    extracted = []
    supported_ext = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.emf', '.wmf', '.tiff'}
    
    try:
        with zipfile.ZipFile(docx_path, 'r') as z:
            for name in z.namelist():
                if name.startswith('word/media/'):
                    ext = Path(name).suffix.lower()
                    if ext in supported_ext:
                        try:
                            data = z.read(name)
                            if data:
                                out_path = output_dir / f"{uuid.uuid4().hex}{ext}"
                                out_path.write_bytes(data)
                                extracted.append(out_path)
                                logger.info(f"zipfile提取DOCX图片: {out_path.name}")
                        except Exception as e:
                            logger.warning(f"提取 {name} 失败: {e}")
    except Exception as e:
        logger.error(f"zipfile提取DOCX图片失败: {e}")
    
    return extracted

async def extract_xlsx_images_safe(xlsx_path: Path) -> list[Path]:
    """
    安全提取Excel图片，处理各种格式问题
    """
    from openpyxl import load_workbook
    
    temp_dir = Path(tempfile.mkdtemp())
    extracted_images = []
    
    try:
        wb = load_workbook(xlsx_path, data_only=True)
        
        # 方法1：通过openpyxl提取（适用于.xlsx）
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # 检查不同版本的openpyxl API
            images_to_process = []
            
            if hasattr(sheet, '_images') and sheet._images:
                images_to_process = sheet._images
            elif hasattr(sheet, 'images') and sheet.images:
                images_to_process = sheet.images
            
            for img_idx, image in enumerate(images_to_process):
                try:
                    # 尝试获取图片数据
                    if hasattr(image, '_data'):
                        image_bytes = image._data
                    elif hasattr(image, 'data'):
                        image_bytes = image.data
                    elif hasattr(image, 'blob'):
                        image_bytes = image.blob
                    else:
                        logger.warning(f"Sheet {sheet_name} 图片 {img_idx}: 无法获取图片数据")
                        continue
                    
                    if not image_bytes:
                        logger.warning(f"Sheet {sheet_name} 图片 {img_idx}: 内容为空")
                        continue
                    
                    # 尝试确定图片格式
                    try:
                        from PIL import Image
                        img = Image.open(io.BytesIO(image_bytes))
                        ext = img.format.lower() if img.format else 'png'
                    except:
                        ext = 'png'
                    
                    img_path = temp_dir / f"xlsx_{sheet_name}_img{img_idx}_{uuid.uuid4().hex}.{ext}"
                    img_path.write_bytes(image_bytes)
                    extracted_images.append(img_path)
                    logger.info(f"成功提取Excel图片: {img_path.name}")
                    
                except Exception as e:
                    logger.warning(f"提取Sheet {sheet_name} 图片 {img_idx} 失败: {type(e).__name__}: {e}")
        
        # 方法2：如果方法1没有提取到图片，尝试zipfile
        if not extracted_images:
            logger.info("尝试使用zipfile直接提取Excel图片...")
            extracted_images = await extract_xlsx_images_via_zip(xlsx_path, temp_dir)
            
    except Exception as e:
        logger.error(f"解析Excel失败: {e}")
        # 尝试zipfile作为备用
        extracted_images = await extract_xlsx_images_via_zip(xlsx_path, temp_dir)
    
    return extracted_images


async def extract_xlsx_images_via_zip(xlsx_path: Path, output_dir: Path) -> list[Path]:
    """
    使用zipfile直接解压Excel中的图片（备用方法）
    """
    extracted = []
    # Excel图片可能在xl/media/或xl/drawings/
    supported_ext = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.emf', '.wmf'}
    
    try:
        with zipfile.ZipFile(xlsx_path, 'r') as z:
            for name in z.namelist():
                # Excel图片通常在 xl/media/ 目录下
                if name.startswith('xl/media/') or name.startswith('xl/drawings/media/'):
                    ext = Path(name).suffix.lower()
                    if ext in supported_ext:
                        try:
                            data = z.read(name)
                            if data:
                                out_path = output_dir / f"{uuid.uuid4().hex}{ext}"
                                out_path.write_bytes(data)
                                extracted.append(out_path)
                                logger.info(f"zipfile提取Excel图片: {out_path.name} (from {name})")
                        except Exception as e:
                            logger.warning(f"提取 {name} 失败: {e}")
    except Exception as e:
        logger.error(f"zipfile提取Excel图片失败: {e}")
    
    return extracted


async def extract_xls_images_safe(xls_path: Path) -> list[Path]:
    """
    提取旧版.xls格式的图片
    注意：.xls是二进制格式，需要特殊处理
    """
    temp_dir = Path(tempfile.mkdtemp())
    extracted_images = []
    
    try:
        # 尝试使用xlrd提取
        import xlrd
        book = xlrd.open_workbook(xls_path, formatting_info=True)
        
        # xlrd不直接支持图片提取，尝试其他方法
        logger.warning(f".xls格式 {xls_path.name} 图片提取受限，尝试备用方法...")
        
        # 尝试使用zipfile（虽然.xls不是zip格式，但某些情况下可能有效）
        # 实际上.xls是OLE格式，不是zip
        logger.info(".xls是二进制OLE格式，无法使用zipfile提取")
        
    except ImportError:
        logger.warning("xlrd未安装，无法处理.xls文件")
    except Exception as e:
        logger.error(f"解析.xls失败: {e}")
    
    # .xls格式通常需要Windows COM或特殊库来处理图片
    # 这里返回空列表，但记录日志
    if not extracted_images:
        logger.warning(f".xls文件 {xls_path.name} 的图片提取需要特殊处理（OLE格式）")
    
    return extracted_images


def is_linked_picture(shape) -> bool:
    """
    检查图片是否是链接图片（非嵌入）
    """
    try:
        from pptx.oxml.ns import qn
        blip = shape._pic.spPr.blipFill.blip
        return blip is not None and blip.get(qn('r:link')) is not None
    except:
        return False


async def extract_images_via_zip(pptx_path: Path, output_dir: Path, existing_hashes: set = None, existing_images: list = None) -> list[Path]:
    """
    使用zipfile直接解压PPT中的图片（绕过python-pptx的bug）
    """
    import hashlib

    if existing_hashes is None:
        existing_hashes = set()
    if existing_images is None:
        existing_images = []
    
    extracted = []
    supported_ext = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            # 只处理 ppt/media/ 目录下的文件
            media_files = [
                name for name in z.namelist() 
                if name.startswith('ppt/media/') and 
                Path(name).suffix.lower() in supported_ext
            ]
            
            logger.info(f"zipfile 发现 {len(media_files)} 个媒体文件")
            
            for name in media_files:
                try:
                    data = z.read(name)
                    if not data or len(data) < 100:  # 过滤空文件和小文件
                        continue
                    
                    # 计算哈希去重
                    img_hash = hashlib.md5(data).hexdigest()
                    if img_hash in existing_hashes:
                        logger.debug(f"zipfile 跳过重复: {name}")
                        continue
                    
                    existing_hashes.add(img_hash)
                    
                    ext = Path(name).suffix.lower()
                    out_path = output_dir / f"zip_{uuid.uuid4().hex}{ext}"
                    out_path.write_bytes(data)
                    extracted.append(out_path)
                    logger.info(f"zipfile提取: {out_path.name} ({len(data)} bytes)")
                    
                except Exception as e:
                    logger.warning(f"提取 {name} 失败: {e}")
    except Exception as e:
        logger.error(f"zipfile提取失败: {e}")
    
    return extracted
 
async def understand_image(
    image: ImagePathOrImage,
    context: str = "",
    detail_level: str = "high",
    language: str = None  # 添加语言参数
) -> str:
    """
    Understand an image with multimodal LLM, extracting both text and meaning.
    
    Args:
        image: Path to image file OR PIL.Image object
        context: Optional context about where this image appears
        detail_level: 'low', 'medium', or 'high' - controls detail level of analysis
    
    Returns:
        Detailed description of image in Markdown format
    """
    from PIL import Image

    # 处理PPT文件 - 提取其中的图片
    if isinstance(image, Path) and image.suffix.lower() == '.pptx':
        logger.info(f"检测到PPT文件，开始提取图片: {image.name}")
        try:
            extracted_images = await extract_ppt_images_safe(image)
            if not extracted_images:
                logger.warning("未能提取到任何图片，可能是链接图片、OLE对象或不支持的格式")
                return "该PPT中没有可提取的嵌入图片（可能包含链接图片、OLE对象或不支持的格式）"
            # 如果有多个图片，使用 understand_multiple_images 处理
            if len(extracted_images) == 1:
                image = extracted_images[0]
            else:
                return await understand_multiple_images(
                    extracted_images, 
                    context=f"{context} (来自PPT: {image.name})"
                )
        except Exception as e:
            logger.error(f"提取PPT图片失败: {e}")
            return f"Error: 无法提取PPT中的图片 - {e}"
    

    # Handle both Path and PIL.Image
    if isinstance(image, Path):
        # From file path
        image_name = image.name
        with open(image, 'rb') as f:
            image_data = f.read()
        
        # Determine MIME type from extension
        mime_type = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp',
            '.tiff': 'image/tiff',
            '.tif': 'image/tiff',
        }.get(image.suffix.lower(), 'image/jpeg')
    else:
        # From PIL.Image object
        image_name = "image.png"
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        image_data = img_byte_arr.getvalue()
        mime_type = 'image/png'
    
    # Encode to base64
    img_base64 = base64.b64encode(image_data).decode('utf-8')
    
    # Get vision model
    model_config = get_vision_model_config()
    # result = await asyncio.to_thread(test_siliconflow_api_config)
    # logger.info(f"SiliconFlow API test result: {result}")
    if not model_config:
        return "Error: No vision-capable model configured"
    
    # Build prompt
    prompt = build_analysis_prompt(detail_level, context, language)
    
    try:
        logger.info(f"Analyzing image: {image_name} (model: {model_config.name})")
        
        # Use direct HTTP request for SiliconFlow models
        if 'siliconflow' in model_config.base_url.lower():
            content = await call_siliconflow_vision_api(
                model_config, img_base64, mime_type, prompt
            )
        else:
            # Use LangChain for other models (GPT-4o, Claude, etc.)
            content = await call_standard_vision_model(
                model_config, img_base64, mime_type, prompt
            )
        
        logger.info(f"Image analysis complete: {len(content)} characters")
        return content
        
    except Exception as e:
        logger.error(f"Failed to understand image {image_name}: {e}")
        return f"Error analyzing image: {e}"
 
 
async def call_siliconflow_vision_api(
    model_config,
    base64_image: str,
    mime_type: str,
    prompt: str
) -> str:
    """Call SiliconFlow vision API directly using requests."""
    import asyncio
    import json
    
    # Build API URL
    base_url = model_config.base_url.rstrip('/')
    api_url = f"{base_url}/chat/completions"
    
    # ========== 确定模型名称 ==========
    model_name = None
    if hasattr(model_config, 'model'):
        model_name = model_config.model
    elif hasattr(model_config, 'name'):
        model_name = model_config.name
    
    # DeepSeek OCR 的完整名称
    if not model_name or 'DeepSeek-OCR' not in model_name:
        model_name = "deepseek-ai/DeepSeek-OCR"
        logger.warning(f"Using default model name: {model_name}")
    # =================================
    
    # ========== 构建 payload ==========
    # 使用传入的 mime_type（image/png 或 image/jpeg）
    image_url = f"data:{mime_type};base64,{base64_image}"
    
    payload = {
        "model": model_name,  # 使用正确的模型名称
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",  # 必须包含text类型
                        "text": prompt   # 使用构建的prompt
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": image_url
                        }  # 直接使用传入的image_url
                    }
                ]
            }
        ],
        "max_tokens": 2048,      # 必须包含max_tokens
        "stream": False
    }
    # ==============================
    
    headers = {
        "Authorization": f"Bearer {model_config.api_key}",
        "Content-Type": "application/json"
    }
    
    # ========== 调试日志 ==========
    logger.info(f"SiliconFlow API URL: {api_url}")
    logger.info(f"Model: {model_name}")
    logger.info(f"Image URL prefix: {image_url[:50]}...")
    logger.info(f"Full payload: {json.dumps(payload, indent=2, ensure_ascii=False)[:1000]}...")
    # ========================
    
    loop = asyncio.get_event_loop()
    
    # ========== 带重试的请求 ==========
    max_retries = 3
    timeout_seconds = 60
    
    for attempt in range(max_retries):
        try:
            def make_request():
                response = requests.post(
                    api_url, 
                    json=payload, 
                    headers=headers, 
                    timeout=timeout_seconds
                )
                logger.info(f"Response status: {response.status_code}")
                if response.status_code != 200:
                    logger.error(f"Response error: {response.text[:500]}")
                # logger.info(f"Response body: {response.text[:1000]}")
                response.raise_for_status()
                return response.json()

            result = await loop.run_in_executor(None, make_request)
            
            # 成功，跳出重试循环
            break
            
        except requests.exceptions.Timeout:
            logger.warning(f"Request timeout (attempt {attempt + 1}/{max_retries})")
            if attempt < max_retries - 1:
                # 使用 asyncio.sleep 在异步上下文中等待
                await asyncio.sleep(3 ** attempt)  # 指数退避：2, 4, 8秒
                continue
            else:
                raise Exception(f"API request timed out after {max_retries} attempts")  # 最后一次重试失败，抛出异常
                
        except requests.exceptions.HTTPError as e:
            # HTTP错误不重试，直接处理
            logger.error(f"HTTP error: {e}")
            raise
        except Exception as e:
            logger.error(f"Unexpected error (attempt {attempt + 1}/{max_retries}): {e}")
            if attempt < max_retries - 1:
                await asyncio.sleep(3 ** attempt)
                continue
            else:
                raise
    # ==================================
    
    try:
        result = await loop.run_in_executor(None, make_request)
        
        
        if "choices" in result and result["choices"]:
            content = result["choices"][0]["message"]["content"]

            # ========== 添加日志看原始内容 ==========
            logger.info(f"Raw response from DeepSeek OCR: {content}")
            # =====================================
            
            # Clean DeepSeek OCR tags
            import re
            content = re.sub(r'<\|ref\|>(.*?)<\|/ref\|>', r'\1', content)
            
            return content
        else:
            logger.error(f"Unexpected response: {json.dumps(result, indent=2)}")
            return "Error: Unexpected API response"
            
    except requests.exceptions.HTTPError as e:
        error_details = ""
        if e.response:
            logger.error(f"HTTP Error Status: {e.response.status_code}")
            logger.error(f"HTTP Error Body: {e.response.text[:1000]}")
            
            try:
                error_json = e.response.json()
                logger.error(f"硅基流动API详细错误: {json.dumps(error_json, indent=2, ensure_ascii=False)}")
                
                # 提取所有可能的错误信息字段
                error_parts = []
                if "error" in error_json:
                    if isinstance(error_json["error"], dict):
                        error_parts.append(f"错误: {error_json['error'].get('message', str(error_json['error']))}")
                    else:
                        error_parts.append(f"错误: {error_json['error']}")
                
                if "message" in error_json:
                    error_parts.append(f"消息: {error_json['message']}")
                
                if "detail" in error_json:
                    error_parts.append(f"详情: {error_json['detail']}")
                
                error_details = " | ".join(error_parts) if error_parts else str(e)
                
            except json.JSONDecodeError:
                error_text = e.response.text[:500]
                logger.error(f"API非JSON错误响应: {error_text}")
                error_details = f"HTTP {e.response.status_code}: {error_text}"
        else:
            error_details = str(e)
        
        return f"Error analyzing image: {error_details}"
 
 
# def build_siliconflow_image_url(
#     base64_image: str,
#     mime_type: str,
#     model_name: str
# ) -> str | dict:
#     """
#     Build image_url for SiliconFlow models.
    
#     DeepSeek OCR expects plain string, others expect nested object.
#     """
#     model_name_lower = model_name.lower()
    
#     if 'deepseek' in model_name_lower and 'ocr' in model_name_lower:
#         # DeepSeek OCR: plain string format
#         return f"data:{mime_type};base64,{base64_image}"
#     else:
#         # Other SiliconFlow models: nested object
#         return {
#             "url": f"data:{mime_type};base64,{base64_image}",
#             "detail": "low"  # SiliconFlow uses 'low' for better performance
#         }
 
def build_siliconflow_image_url(
    base64_image: str,
    mime_type: str,
    model_name: str
) -> str | dict:
    """
    Build image_url for SiliconFlow models.
    
    DeepSeek OCR: 使用 image/png 或 image/jpeg
    """
    # 直接使用传入的 mime_type
    return f"data:{mime_type};base64,{base64_image}"
 
async def call_standard_vision_model(
    model_config,
    base64_image: str,
    mime_type: str,
    prompt: str
) -> str:
    """
    Call standard vision models (GPT-4o, Claude, etc.) via LangChain.
    """
    from src.models.factory import create_chat_model
    
    try:
        model = create_chat_model(name=model_config.name)
        
        # Build message
        content = [
            {
                "type": "text",
                "text": prompt
            },
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{mime_type};base64,{base64_image}",
                    "detail": "high"
                }
            }
        ]
        
        message = HumanMessage(content=content)
        response = await model.ainvoke([message])
        
        # Extract content
        if isinstance(response.content, list):
            content = '\n\n'.join(
                item.text if hasattr(item, 'text') else str(item)
                for item in response.content
            )
        else:
            content = str(response.content)
        
        return content
        
    except Exception as e:
        logger.error(f"Standard vision model failed: {e}")
        return f"Error: {e}"
 
 
# def build_analysis_prompt(detail_level: str, context: str) -> str:
#     """Build analysis prompt based on detail level."""
#     if detail_level == "high":
#         return f"""Analyze this image in detail and provide:
 
# 1. **Type of Image**: (photo, chart, graph, diagram, screenshot, document scan, etc.)
 
# 2. **Content Summary**: Brief overview of what the image shows
 
# 3. **Text Content**: Extract ALL visible text, word-for-word. Preserve formatting and structure as much as possible.
 
# 4. **Key Elements**: 
#    - For charts/graphs: List axes, labels, data points, trends
#    - For diagrams: Explain components and relationships
#    - For documents: Describe layout, sections, formatting
#    - For photos: Describe main subjects, actions, setting
 
# 5. **Data Extraction** (if applicable):
#    - Tables: Convert to Markdown table format
#    - Charts: Describe data and insights
#    - Forms: List all fields and their values
 
# 6. **Context & Meaning**: Explain the significance of this image in the document.
 
# {f"Additional Context: {context}" if context else ""}
 
# Provide the response in Markdown format with proper formatting. Be thorough and precise."""
    
#     elif detail_level == "medium":
#         return f"""Analyze this image and provide:
 
# 1. **Image Type**: What kind of image is this?
 
# 2. **Text Content**: Extract all visible text accurately.
 
# 3. **Main Content**: Brief description of what the image shows.
 
# 4. **Key Information**: Highlight important data, insights, or elements.
 
# {f"Context: {context}" if context else ""}
 
# Provide in Markdown format."""
    
#     else:  # low
#         return f"""Describe this image in detail. Extract any text content. Provide a clear, concise description."""

def build_analysis_prompt(detail_level: str, context: str, language: str = None) -> str:
    """
    Build analysis prompt based on detail level and language.
    
    Args:
        detail_level: 'low', 'medium', or 'high'
        context: Optional context about the image
        language: 'zh' for Chinese, 'en' for English, None for auto-detect from config
    
    Returns:
        Prompt string in the specified language
    """
    # 如果未指定语言，使用配置
    if language is None:
        language = get_language()
    
    # 定义双语提示词模板
    PROMPTS = {
        "zh": {
            "high": """请详细分析这张图片并提供：

1. **图片类型**：(照片、图表、图形、示意图、截图、文档扫描等)

2. **内容摘要**：简要概述图片显示的内容

3. **文本内容**：提取所有可见文本，逐字提取。尽可能保留格式和结构。

4. **关键元素**：
   - 对于图表/图形：列出坐标轴、标签、数据点、趋势
   - 对于示意图：解释组件和关系
   - 对于文档：描述布局、章节、格式
   - 对于照片：描述主要主体、动作、场景

5. **数据提取**（如适用）：
   - 表格：转换为Markdown表格格式
   - 图表：描述数据和洞察
   - 表单：列出所有字段及其值

6. **上下文与意义**：解释这张图片在文档中的重要性。

{f"额外上下文：{context}" if context else ""}

请用Markdown格式提供回复，格式规范。请详尽而准确。""",

            "medium": """分析这张图片并提供：

1. **图片类型**：这是什么类型的图片？

2. **文本内容**：准确提取所有可见文本。

3. **主要内容**：简要描述图片显示的内容。

4. **关键信息**：突出重要的数据、洞察或元素。

{f"上下文：{context}" if context else ""}

请用Markdown格式提供回复。""",

            "low": """详细描述这张图片。提取任何文本内容。提供清晰、简洁的描述。"""
        },
        "en": {
            "high": """Analyze this image in detail and provide:

1. **Type of Image**: (photo, chart, graph, diagram, screenshot, document scan, etc.)

2. **Content Summary**: Brief overview of what the image shows

3. **Text Content**: Extract ALL visible text, word-for-word. Preserve formatting and structure as much as possible.

4. **Key Elements**: 
   - For charts/graphs: List axes, labels, data points, trends
   - For diagrams: Explain components and relationships
   - For documents: Describe layout, sections, formatting
   - For photos: Describe main subjects, actions, setting

5. **Data Extraction** (if applicable):
   - Tables: Convert to Markdown table format
   - Charts: Describe data and insights
   - Forms: List all fields and their values

6. **Context & Meaning**: Explain the significance of this image in the document.

{f"Additional Context: {context}" if context else ""}

Provide the response in Markdown format with proper formatting. Be thorough and precise.""",

            "medium": """Analyze this image and provide:

1. **Image Type**: What kind of image is this?

2. **Text Content**: Extract all visible text accurately.

3. **Main Content**: Brief description of what the image shows.

4. **Key Information**: Highlight important data, insights, or elements.

{f"Context: {context}" if context else ""}

Provide in Markdown format.""",

            "low": """Describe this image in detail. Extract any text content. Provide a clear, concise description."""
        }
    }
    
    # 验证语言参数
    if language not in PROMPTS:
        logger.warning(f"Unsupported language: {language}, falling back to {DEFAULT_LANGUAGE}")
        language = DEFAULT_LANGUAGE
    
    # 验证 detail_level
    if detail_level not in ["high", "medium", "low"]:
        detail_level = "high"
    
    return PROMPTS[language][detail_level]

 
async def understand_multiple_images(
    images: list[ImagePathOrImage],
    context: str = "",
    detail_level: str = "high",
    language: str = None  # 添加语言参数
) -> str:
    """
    Understand multiple images with context awareness.
    
    Args:
        images: List of image paths OR PIL.Image objects
        context: Context about where these images appear
    
    Returns:
        Combined analysis in Markdown format
    """
    analyses = []
    
    async with temp_image_manager([img for img in images if isinstance(img, Path)]):
        for i, img in enumerate(images, 1):
            # Get image name
            if isinstance(img, Path):
                img_name = img.name
            else:
                img_name = f"image_{i}.png"

            logger.info(f"Processing image {i}/{len(images)}: {img_name}")

            try:
                analysis = await understand_image(
                    img,
                    context=f"{context} (Image {i} of {len(images)})",
                    detail_level=detail_level,
                    language=language  # 传递语言参数
                )
                # 根据语言选择标题
                if language == "en" or (language is None and get_language() == "en"):
                    analyses.append(f"### Image {i}: {img_name}\n\n{analysis}")
                else:
                    analyses.append(f"### 图片 {i}: {img_name}\n\n{analysis}")
            except Exception as e:
                logger.error(f"Failed to analyze image {img_name}: {e}")
                analyses.append(f"### Image {i}: {img_name}\n\nError analyzing image: {e}")

            # Small delay to avoid rate limiting
            await asyncio.sleep(0.5)

    return '\n\n---\n\n'.join(analyses)
 
 
def get_vision_model_config():
    """
    Get vision model configuration.
    
    Returns:
        ModelConfig object or None
    """
    try:
        from src.config.app_config import get_app_config
        
        config = get_app_config()
        
        # Priority order: SiliconFlow models first
        preferred_models = [
            'siliconflow-deepseek-ocr',
            'siliconflow-paddle-ocr',
            'qwen-vl',
            'gpt-4o',
            'gpt-4-turbo',
            'claude-3.5-sonnet',
            'claude-3-opus',
        ]
        
        for model_name in preferred_models:
            for model_config in config.models:
                if model_name.lower() in model_config.name.lower():
                    logger.info(f"Using vision model: {model_config.name}")
                    return model_config
        
        logger.error(
            "No vision-capable model found. "
            "Please configure: siliconflow-deepseek-ocr, PaddleOCR, GPT-4o, etc."
        )
        return None
        
    except Exception as e:
        logger.error(f"Failed to get vision model config: {e}")
        return None
 
 
def sanitize_analysis_output(content: str, max_length: int = 200000) -> str:
    """Sanitize LLM analysis output to prevent encoding and size issues."""
    if not content:
        return ""
    
    import re
    
    # Remove control characters
    content = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', content)
    
    # Fix Unicode issues
    content = content.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
    
    # Normalize whitespace
    content = re.sub(r'\n{3,}', '\n\n', content)
    
    # Truncate if needed
    if len(content) > max_length:
        content = content[:max_length] + (
            f"\n\n[Content truncated from {len(content)} to {max_length} characters]"
        )
    
    return content.strip()

# def test_siliconflow_api_config():
#     """测试硅基流动API配置和连通性"""
#     import requests
    
#     # 获取配置
#     model_config = get_vision_model_config()
#     if not model_config:
#         return "❌ 未找到模型配置"
    
#     # 构建最小测试请求
#     test_payload = {
#         "model": "deepseek-ocr",
#         "messages": [
#             {
#                 "role": "user",
#                 "content": [
#                     {"type": "text", "text": "Hello, can you read this?"}
#                 ]
#             }
#         ],
#         "max_tokens": 100
#     }
    
#     headers = {
#         "Authorization": f"Bearer {model_config.api_key}",
#         "Content-Type": "application/json"
#     }
    
#     api_url = f"{model_config.base_url.rstrip('/')}/chat/completions"
    
#     try:
#         response = requests.post(api_url, json=test_payload, headers=headers, timeout=10)
        
#         if response.status_code == 200:
#             return "✅ API连通性测试成功"
#         else:
#             return f"❌ API测试失败: {response.status_code}\n{response.text[:200]}"
            
#     except Exception as e:
#         return f"❌ API测试异常: {e}"

def test_siliconflow_api_config():
    """测试硅基流动API配置和连通性"""
    import requests
    
    model_config = get_vision_model_config()
    if not model_config:
        return "❌ 未找到模型配置"
    
    # 测试1：纯文本请求（验证API密钥和连通性）
    test_payload = {
        "model": "deepseek-ai/DeepSeek-OCR",
        "messages": [
            {
                "role": "user",
                "content": "Hello, are you working?"
            }
        ],
        "max_tokens": 100,
        "temperature": 0.1
    }
    
    headers = {
        "Authorization": f"Bearer {model_config.api_key}",
        "Content-Type": "application/json"
    }
    
    api_url = f"{model_config.base_url.rstrip('/')}/chat/completions"
    
    try:
        response = requests.post(api_url, json=test_payload, headers=headers, timeout=10)
        print(f"Test response: {response.status_code}")
        print(f"Test body: {response.text[:500]}")
        
        if response.status_code == 200:
            return "✅ API连通性测试成功"
        else:
            return f"❌ API测试失败: {response.status_code}\n{response.text[:500]}"
    except Exception as e:
        return f"❌ API测试异常: {e}"