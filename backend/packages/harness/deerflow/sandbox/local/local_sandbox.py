import os
import shutil
import subprocess
from pathlib import Path

import hashlib
import io
import json
import os
import re
from typing import Any

import chardet

from src.sandbox.local.list_dir import list_dir
from src.sandbox.sandbox import Sandbox


class LocalSandbox(Sandbox):
    def __init__(self, id: str, path_mappings: dict[str, str] | None = None):
        """
        Initialize local sandbox with optional path mappings.

        Args:
            id: Sandbox identifier
            path_mappings: Dictionary mapping container paths to local paths
                          Example: {"/mnt/skills": "/absolute/path/to/skills"}
        """
        super().__init__(id)
        self.path_mappings = path_mappings or {}

    def _resolve_path(self, path: str) -> str:
        """
        Resolve container path to actual local path using mappings.

        Args:
            path: Path that might be a container path

        Returns:
            Resolved local path
        """
        path_str = str(path)

        # Try each mapping (longest prefix first for more specific matches)
        for container_path, local_path in sorted(self.path_mappings.items(), key=lambda x: len(x[0]), reverse=True):
            if path_str.startswith(container_path):
                # Replace the container path prefix with local path
                relative = path_str[len(container_path) :].lstrip("/")
                resolved = str(Path(local_path) / relative) if relative else local_path
                return resolved

        # No mapping found, return original path
        return path_str

    def _reverse_resolve_path(self, path: str) -> str:
        """
        Reverse resolve local path back to container path using mappings.

        Args:
            path: Local path that might need to be mapped to container path

        Returns:
            Container path if mapping exists, otherwise original path
        """
        path_str = str(Path(path).resolve())

        # Try each mapping (longest local path first for more specific matches)
        for container_path, local_path in sorted(self.path_mappings.items(), key=lambda x: len(x[1]), reverse=True):
            local_path_resolved = str(Path(local_path).resolve())
            if path_str.startswith(local_path_resolved):
                # Replace the local path prefix with container path
                relative = path_str[len(local_path_resolved) :].lstrip("/")
                resolved = f"{container_path}/{relative}" if relative else container_path
                return resolved

        # No mapping found, return original path
        return path_str

    def _reverse_resolve_paths_in_output(self, output: str) -> str:
        """
        Reverse resolve local paths back to container paths in output string.

        Args:
            output: Output string that may contain local paths

        Returns:
            Output with local paths resolved to container paths
        """
        import re

        # Sort mappings by local path length (longest first) for correct prefix matching
        sorted_mappings = sorted(self.path_mappings.items(), key=lambda x: len(x[1]), reverse=True)

        if not sorted_mappings:
            return output

        # Create pattern that matches absolute paths
        # Match paths like /Users/... or other absolute paths
        result = output
        for container_path, local_path in sorted_mappings:
            local_path_resolved = str(Path(local_path).resolve())
            # Escape the local path for use in regex
            escaped_local = re.escape(local_path_resolved)
            # Match the local path followed by optional path components
            pattern = re.compile(escaped_local + r"(?:/[^\s\"';&|<>()]*)?")

            def replace_match(match: re.Match) -> str:
                matched_path = match.group(0)
                return self._reverse_resolve_path(matched_path)

            result = pattern.sub(replace_match, result)

        return result

    def _resolve_paths_in_command(self, command: str) -> str:
        """
        Resolve container paths to local paths in a command string.

        Args:
            command: Command string that may contain container paths

        Returns:
            Command with container paths resolved to local paths
        """
        import re

        # Sort mappings by length (longest first) for correct prefix matching
        sorted_mappings = sorted(self.path_mappings.items(), key=lambda x: len(x[0]), reverse=True)

        # Build regex pattern to match all container paths
        # Match container path followed by optional path components
        if not sorted_mappings:
            return command

        # Create pattern that matches any of the container paths
        patterns = [re.escape(container_path) + r"(?:/[^\s\"';&|<>()]*)??" for container_path, _ in sorted_mappings]
        pattern = re.compile("|".join(f"({p})" for p in patterns))

        def replace_match(match: re.Match) -> str:
            matched_path = match.group(0)
            return self._resolve_path(matched_path)

        return pattern.sub(replace_match, command)

    @staticmethod
    def _get_shell() -> str:
        """Detect available shell executable with fallback.

        Returns the first available shell in order of preference:
        /bin/zsh → /bin/bash → /bin/sh → first `sh` found on PATH.
        Raises a RuntimeError if no suitable shell is found.
        """
        import platform
        import shutil

        # Windows support
        if platform.system() == "Windows":
            # Try PowerShell first, then cmd
            if shutil.which("powershell.exe"):
                return "powershell.exe"
            if shutil.which("cmd.exe"):
                return "cmd.exe"
            raise RuntimeError("No suitable shell found on Windows. Tried powershell.exe and cmd.exe")
        
        for shell in ("/bin/zsh", "/bin/bash", "/bin/sh"):
            if os.path.isfile(shell) and os.access(shell, os.X_OK):
                return shell
        shell_from_path = shutil.which("sh")
        if shell_from_path is not None:
            return shell_from_path
        raise RuntimeError("No suitable shell executable found. Tried /bin/zsh, /bin/bash, /bin/sh, and `sh` on PATH.")

    def execute_command(self, command: str) -> str:
        # Resolve container paths in command before execution
        resolved_command = self._resolve_paths_in_command(command)

        result = subprocess.run(
            resolved_command,
            executable=self._get_shell(),
            shell=True,
            capture_output=True,
            text=True,
            timeout=600,
        )
        output = result.stdout
        if result.stderr:
            output += f"\nStd Error:\n{result.stderr}" if output else result.stderr
        if result.returncode != 0:
            output += f"\nExit Code: {result.returncode}"

        final_output = output if output else "(no output)"
        # Reverse resolve local paths back to container paths in output
        return self._reverse_resolve_paths_in_output(final_output)

    def list_dir(self, path: str, max_depth=2) -> list[str]:
        resolved_path = self._resolve_path(path)
        entries = list_dir(resolved_path, max_depth)
        # Reverse resolve local paths back to container paths in output
        return [self._reverse_resolve_paths_in_output(entry) for entry in entries]

    # def read_file(self, path: str) -> str:
    #     resolved_path = self._resolve_path(path)
    #     try:
    #         with open(resolved_path, encoding="utf-8") as f:
    #             return f.read()
    #     except OSError as e:
    #         # Re-raise with the original path for clearer error messages, hiding internal resolved paths
    #         raise type(e)(e.errno, e.strerror, path) from None

    # def write_file(self, path: str, content: str, append: bool = False) -> None:
    #     resolved_path = self._resolve_path(path)
    #     try:
    #         dir_path = os.path.dirname(resolved_path)
    #         if dir_path:
    #             os.makedirs(dir_path, exist_ok=True)
    #         mode = "a" if append else "w"
    #         with open(resolved_path, mode, encoding="utf-8") as f:
    #             f.write(content)
    #     except OSError as e:
    #         # Re-raise with the original path for clearer error messages, hiding internal resolved paths
    #         raise type(e)(e.errno, e.strerror, path) from None

    # def update_file(self, path: str, content: bytes) -> None:
    #     resolved_path = self._resolve_path(path)
    #     try:
    #         dir_path = os.path.dirname(resolved_path)
    #         if dir_path:
    #             os.makedirs(dir_path, exist_ok=True)
    #         with open(resolved_path, "wb") as f:
    #             f.write(content)
    #     except OSError as e:
    #         # Re-raise with the original path for clearer error messages, hiding internal resolved paths
    #         raise type(e)(e.errno, e.strerror, path) from None
    
    # def read_file(self, path: str) -> str:
    #     resolved_path = self._resolve_path(path)
    #     with open(resolved_path) as f:
    #         return f.read()
    
    def _sanitize_markdown_content(self, content: str, max_length: int = 100000) -> str:
        """Sanitize markdown content to prevent encoding errors."""
        if not content:
            return ""

        # Remove control characters except newlines and tabs
        content = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', content)

        # Fix Unicode issues
        content = content.encode('utf-8', errors='replace').decode('utf-8', errors='replace')

        # Normalize whitespace
        content = re.sub(r'[ \t]+', ' ', content)
        content = re.sub(r'\n{3,}', '\n\n', content)

        # Truncate if too large
        if len(content) > max_length:
            content = content[:max_length] + (
                f"\n\n[Content truncated from {len(content)} to {max_length} characters]"
            )

        return content.strip()

    def _is_binary_file(self, file_path: str | Path) -> bool:
        """Check if a file is binary by reading first 1024 bytes."""
        try:
            with open(file_path, 'rb') as f:
                chunk = f.read(1024)
                # Check for null bytes
                if b'\x00' in chunk:
                    return True
                # Check ratio of non-text bytes
                text_chars = bytearray(
                    {7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7f}
                )
                non_text = sum(1 for byte in chunk if byte not in text_chars)
                return non_text / len(chunk) > 0.3
        except Exception:
            return True

    def _load_file_metadata(self, file_path: Path) -> dict:
        """Load metadata for a file."""
        if isinstance(file_path, str):
            file_path = Path(file_path)
        meta_path = file_path.with_suffix('.meta.json')
        if meta_path.exists():
            try:
                with open(meta_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                pass
        return {}

    # ============ 修改：read_file 方法 ============

    async def read_file(
            self,
            path: str,
            start_line: int = None,
            end_line: int = None,
            max_size: int = 100000,
        ) -> str:
        """
        Read file with intelligent handling of images and scanned content.

        Strategy:
        - Plain text/markdown: Read directly
        - Scanned PDF: OCR with LLM
        - PDF with images: Extract text + analyze images
        - Office with images: Extract content + analyze images
        - Images: Analyze with LLM
        """

        resolved_path_str = self._resolve_path(path)
        resolved_path = Path(resolved_path_str)

        if not os.path.exists(resolved_path):
            return f"Error: File not found: {resolved_path}"

        # Check file size
        try:
            file_size = os.path.getsize(resolved_path)
            if file_size > max_size:
                return (
                    f"Error: File is too large ({file_size} bytes). "
                    f"Maximum allowed: {max_size} bytes."
                )
        except Exception as e:
            return f"Error checking file size: {e}"

        # Get file extension
        # _, ext = os.path.splitext(resolved_path)
        # ext = ext.lower()
        ext = resolved_path.suffix.lower()

        # Load metadata
        metadata = self._load_file_metadata(resolved_path)

        # Handle different file types
        if ext == '.md':
            return await self._handle_markdown_file(resolved_path, metadata)
        elif ext == '.pdf':
            return await self._handle_pdf_file(resolved_path, metadata)
        elif ext in ['.docx', '.pptx', '.xlsx']:
            return await self._handle_office_file(resolved_path, metadata)
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif']:
            return await self._handle_image_file(resolved_path)
        else:
            return await self._handle_text_file(resolved_path, start_line, end_line)

    # ============ 文件处理方法 ============

    async def _handle_markdown_file(self, md_path: Path, metadata: dict) -> str:
        """Handle markdown files, including placeholders."""
        if isinstance(md_path, str):
            md_path = Path(md_path)
        
        # Check if this is a placeholder
        file_type = metadata.get('type', '')

        if 'scanned_pdf' in file_type:
            # This is a scanned PDF placeholder
            return await self._process_scanned_pdf(md_path.with_suffix('.pdf'))

        elif 'with_images' in file_type:
            # This is a placeholder for file with images
            # Try to find original file
            original_file = md_path.with_suffix('')
            if not original_file.exists():
                # Try common extensions
                for ext in ['.pdf', '.docx', '.pptx']:
                    candidate = md_path.with_suffix(ext)
                    if candidate.exists():
                        original_file = candidate
                        break

            if original_file.exists():
                return await self._process_file_with_images(original_file, metadata)

        # Regular markdown file
        try:
            with open(md_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return self._sanitize_markdown_content(content)
        except Exception as e:
            return f"Error reading markdown file: {e}"

    async def _handle_pdf_file(self, pdf_path: Path, metadata: dict) -> str:
        """Handle PDF files."""
        needs_analysis = metadata.get('needs_image_analysis', False)

        if needs_analysis:
            return await self._process_file_with_images(pdf_path, metadata)
        else:
            # Just extract text
            return self._extract_text_from_pdf(pdf_path)

    async def _handle_office_file(self, office_path: Path, metadata: dict) -> str:
        """Handle Office files."""
        needs_analysis = metadata.get('needs_image_analysis', False)

        if needs_analysis:
            return await self._process_file_with_images(office_path, metadata)
        else:
            # Use markitdown for conversion
            return await self._convert_office_with_markitdown(office_path)

    async def _handle_image_file(self, img_path: Path) -> str:
        """Handle image files with LLM analysis."""
        try:
            # Import here to avoid circular imports
            from src.tools.builtins.image_understanding import (
                understand_image,
                sanitize_analysis_output,
            )

            analysis = await understand_image(img_path, detail_level="high")
            return sanitize_analysis_output(analysis)
        except Exception as e:
            return f"Error analyzing image: {e}"

    async def _handle_text_file(
            self,
            file_path: Path,
            start_line: int = None,
            end_line: int = None,
            max_size: int = 100000,
    ) -> str:
        """Handle text files with encoding detection."""
        # Get file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        # Office documents that weren't converted
        office_docs = {
            '.docx': 'Word document',
            '.xlsx': 'Excel spreadsheet',
            '.pptx': 'PowerPoint presentation',
            '.doc': 'Word document',
            '.xls': 'Excel spreadsheet',
            '.ppt': 'PowerPoint presentation',
        }

        if ext in office_docs:
            return (
                f"[Binary file: {office_docs[ext]}. "
                f"This file should have been converted to Markdown during upload. "
                f"Please upload it again or check if the conversion failed.]"
            )

        # PDF
        if ext == '.pdf':
            return (
                f"[Binary file: PDF document. "
                f"This file should have been converted to Markdown during upload. "
                f"Please upload it again or check if the conversion failed.]"
            )

        # Images
        image_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif'}
        if ext in image_exts:
            return (
                f"[Binary file: Image. "
                f"This file should be processed with OCR or image analysis.]"
            )

        # Archives
        archive_exts = {'.zip', '.tar', '.gz', '.rar', '.7z'}
        if ext in archive_exts:
            return (
                f"[Binary file: Archive. "
                f"Extract files first before reading.]"
            )

        # Check if binary
        if self._is_binary_file(file_path):
            return (
                f"[Binary file: {ext}. "
                f"Cannot read as text. Use appropriate tool.]"
            )

        # Read text files with encoding detection
        encodings_to_try = [
            'utf-8-sig',
            'utf-8',
            'gb18030',
            'gbk',
            'gb2312',
            'cp936',
            'latin-1',
        ]

        for encoding in encodings_to_try:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    if start_line is not None and end_line is not None:
                        lines = []
                        for i, line in enumerate(f, 1):
                            if i >= start_line and i <= end_line:
                                lines.append(line)
                            elif i > end_line:
                                break
                        content = ''.join(lines)
                    else:
                        content = f.read()

                    return self._sanitize_markdown_content(content)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                return f"Error reading file with {encoding}: {str(e)}"

        # Fallback: use chardet
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(min(50000, max_size))

            result = chardet.detect(raw_data)
            detected_encoding = result['encoding']

            if detected_encoding is None:
                detected_encoding = 'utf-8'
            elif detected_encoding.lower() in ['gb2312', 'gbk']:
                detected_encoding = 'gb18030'

            with open(file_path, 'r', encoding=detected_encoding, errors='replace') as f:
                content = f.read()

            return self._sanitize_markdown_content(content)
        except Exception as e:
            # Last resort: ignore errors
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return self._sanitize_markdown_content(content)
            except Exception as e2:
                return f"Error reading file {file_path}: {str(e2)}"

    # ============ 图片和扫描处理方法 ============

    async def _process_scanned_pdf(self, pdf_path: Path) -> str:
        """Process a scanned PDF with full OCR."""
        logger = __import__('logging').getLogger(__name__)
        logger.info(f"Processing scanned PDF: {pdf_path.name}")

        try:
            from pdf2image import convert_from_path
            from src.tools.builtins.image_understanding import (
                understand_multiple_images,
                sanitize_analysis_output,
            )

            # Convert PDF to images
            logger.info("Converting PDF to images...")
            images = convert_from_path(
                str(pdf_path),
                dpi=200,
                first_page=1,
                last_page=10  # Limit to 10 pages
            )

            if not images:
                return "Error: Failed to convert PDF to images"

            # Analyze each image
            analyses = await understand_multiple_images(
                images,
                context=f"Scanned PDF document: {pdf_path.name}"
            )

            return sanitize_analysis_output(analyses, max_length=500000)

        except ImportError:
            return "Error: pdf2image not installed. Install with: pip install pdf2image"
        except Exception as e:
            logger.error(f"Failed to process scanned PDF: {e}")
            return f"Error: {e}"

    async def _process_file_with_images(self, file_path: Path, metadata: dict) -> str:
        """Process a file that contains images."""
        logger = __import__('logging').getLogger(__name__)
        logger.info(f"Processing file with images: {file_path.name}")

        ext = file_path.suffix.lower()

        if ext == '.pdf':
            return await self._process_pdf_with_images(file_path, metadata)
        elif ext == '.docx':
            return await self._process_docx_with_images(file_path, metadata)
        elif ext == '.pptx':
            return await self._process_pptx_with_images(file_path, metadata)
        else:
            return await self._convert_office_with_markitdown(file_path)

    async def _process_pdf_with_images(self, pdf_path: Path, metadata: dict) -> str:
        """Process PDF with images: extract text + analyze images."""
        try:
            import PyPDF2
            from src.tools.builtins.image_understanding import understand_multiple_images
            from pdf2image import convert_from_path

            logger = __import__('logging').getLogger(__name__)

            # Extract text from text layer
            logger.info("Extracting text from PDF...")
            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n\n{text}")

                text_content = '\n\n'.join(text_parts)

            # Get pages with images
            pages_with_images = metadata.get('pages_with_images', [])

            if not pages_with_images:
                # No images to analyze
                return self._sanitize_markdown_content(text_content)

            # Extract and analyze images
            logger.info(f"Analyzing {len(pages_with_images)} images...")
            images_to_analyze = []

            all_pages = convert_from_path(str(pdf_path), dpi=200)

            for page_num in pages_with_images:
                if page_num <= len(all_pages):
                    images_to_analyze.append(all_pages[page_num - 1])

            if images_to_analyze:
                image_analysis = await understand_multiple_images(
                    images_to_analyze,
                    context=f"Images from PDF: {pdf_path.name}"
                )
            else:
                image_analysis = ""

            # Combine text and image analysis
            result = f"""# 📄 Document: {pdf_path.name}

            ## Text Content

            {text_content}

            ## Image Analysis

            {image_analysis}
            """

            return self._sanitize_markdown_content(result, max_length=500000)

        except Exception as e:
            logger = __import__('logging').getLogger(__name__)
            logger.error(f"Failed to process PDF with images: {e}")
            return f"Error: {e}"

    async def _process_docx_with_images(self, docx_path: Path, metadata: dict) -> str:
        """Process DOCX with images: extract text + analyze images."""
        try:
            from docx import Document
            from src.tools.builtins.image_understanding import understand_multiple_images
            from PIL import Image

            logger = __import__('logging').getLogger(__name__)
            logger.info("Processing DOCX with images...")

            # Extract text
            doc = Document(docx_path)
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            text_content = '\n\n'.join(text_parts)

            # Extract images
            image_data_list = []
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image = rel.target_part
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    image_data_list.append(img)

            # Analyze images
            image_analysis = ""
            if image_data_list:
                logger.info(f"Analyzing {len(image_data_list)} images...")
                image_analysis = await understand_multiple_images(
                    image_data_list,
                    context=f"Images from document: {docx_path.name}"
                )

            # Combine
            result = f"""# 📄 Document: {docx_path.name}

            ## Text Content

            {text_content}

            ## Image Analysis

            {image_analysis}
            """

            return self._sanitize_markdown_content(result, max_length=500000)

        except Exception as e:
            logger = __import__('logging').getLogger(__name__)
            logger.error(f"Failed to process DOCX: {e}")
            return f"Error: {e}"

    async def _process_pptx_with_images(self, pptx_path: Path, metadata: dict) -> str:
        """Process PPTX with images: extract text + analyze images."""
        try:
            from pptx import Presentation
            from src.tools.builtins.image_understanding import understand_multiple_images
            from PIL import Image

            logger = __import__('logging').getLogger(__name__)
            logger.info("Processing PPTX with images...")

            # Extract text
            prs = Presentation(pptx_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"### Slide {slide_num}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                text_parts.append(slide_text)

            text_content = '\n\n'.join(text_parts)

            # Extract images
            image_data_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        image_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        image_data_list.append(img)

            # Analyze images
            image_analysis = ""
            if image_data_list:
                logger.info(f"Analyzing {len(image_data_list)} images...")
                image_analysis = await understand_multiple_images(
                    image_data_list,
                    context=f"Images from presentation: {pptx_path.name}"
                )

            # Combine
            result = f"""# 📊 Presentation: {pptx_path.name}

            ## Slide Content

            {text_content}

            ## Image Analysis

            {image_analysis}
            """

            return self._sanitize_markdown_content(result, max_length=500000)

        except Exception as e:
            logger = __import__('logging').getLogger(__name__)
            logger.error(f"Failed to process PPTX: {e}")
            return f"Error: {e}"

    async def _convert_office_with_markitdown(self, file_path: Path) -> str:
        """Convert Office file using markitdown."""
        try:
            from markitdown import MarkItDown

            md = MarkItDown(enable_plugins=False)
            result = md.convert(str(file_path))

            if result.text_content:
                return self._sanitize_markdown_content(result.text_content)
            else:
                return "Error: Failed to convert file"
        except Exception as e:
            return f"Error: {e}"

    def _extract_text_from_pdf(self, pdf_path: Path) -> str:
        """Extract text from PDF using PyPDF2."""
        try:
            import PyPDF2

            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)

                return self._sanitize_markdown_content('\n\n'.join(text_parts))
        except Exception as e:
            return f"Error: {e}"

    # def write_file(self, path: str, content: str, append: bool = False) -> None:
    #     resolved_path = self._resolve_path(path)
    #     dir_path = os.path.dirname(resolved_path)
    #     if dir_path:
    #         os.makedirs(dir_path, exist_ok=True)
    #     mode = "a" if append else "w"
    #     with open(resolved_path, mode) as f:
    #         f.write(content)

    def write_file(
            self,
            path: str,
            content: str,
            append: bool = False,
            max_size: int = 1000000,  # 默认 1MB
        ) -> str:
        """
        Write content to a file with proper encoding and validation.

        Args:
            path: File path to write to
            content: Content to write
            append: If True, append to existing file; otherwise overwrite
            max_size: Maximum file size in bytes (default: 1MB)

        Returns:
            Success message or error message
        """
        try:
            resolved_path = self._resolve_path(path)

            # Create directory if needed
            dir_path = os.path.dirname(resolved_path)
            if dir_path:
                os.makedirs(dir_path, exist_ok=True)

            # Sanitize content before writing
            sanitized_content = self._sanitize_write_content(content, max_size)

            # Check if appending would exceed limit
            if append and os.path.exists(resolved_path):
                existing_size = os.path.getsize(resolved_path)
                if existing_size + len(sanitized_content) > max_size:
                    return (
                        f"Error: Cannot append. File would exceed size limit "
                        f"({existing_size + len(sanitized_content)} > {max_size} bytes)"
                    )

            # Determine mode
            mode = "a" if append else "w"

            # Write with UTF-8 encoding
            with open(resolved_path, mode, encoding='utf-8') as f:
                f.write(sanitized_content)

            action = "appended to" if append else "written to"
            bytes_written = len(sanitized_content)

            return (
                f"Success: Content {action} {path} "
                f"({bytes_written} characters)"
            )

        except PermissionError:
            return f"Error: Permission denied when writing to {path}"
        except IsADirectoryError:
            return f"Error: {path} is a directory, not a file"
        except Exception as e:
            return f"Error writing to {path}: {str(e)}"

    def _sanitize_write_content(self, content: str, max_size: int) -> str:
        """
        Sanitize content before writing to prevent encoding and size issues.

        Args:
            content: Raw content
            max_size: Maximum allowed size

        Returns:
            Sanitized content
        """
        if not content:
            return ""

        # Ensure content is string
        if not isinstance(content, str):
            try:
                content = str(content)
            except Exception as e:
                raise ValueError(f"Cannot convert content to string: {e}")

        # Remove control characters except common ones
        content = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', content)

        # Fix Unicode issues
        content = content.encode('utf-8', errors='replace').decode('utf-8', errors='replace')

        # Normalize whitespace
        content = re.sub(r'[ \t]+', ' ', content)
        content = re.sub(r'\n{3,}', '\n\n', content)

        # Check size and truncate if needed
        if len(content) > max_size:
            truncated = content[:max_size]
            truncated += (
                f"\n\n[Content truncated from {len(content)} to {max_size} characters. "
                f"This is to prevent file size issues.]"
            )
            return truncated

        return content.strip()


    # def update_file(self, path: str, content: bytes) -> None:
    #     resolved_path = self._resolve_path(path)
    #     dir_path = os.path.dirname(resolved_path)
    #     if dir_path:
    #         os.makedirs(dir_path, exist_ok=True)
    #     with open(resolved_path, "wb") as f:
    #         f.write(content)

    def update_file(
            self,
            path: str,
            content: bytes,
            max_size: int = 524288000,  # 默认 500MB
    ) -> str:
        """
        Write binary content to a file.

        Args:
            path: File path to write to
            content: Binary content to write
            max_size: Maximum file size in bytes (default: 10MB)

        Returns:
            Success message or error message

        Note:
            This method handles binary files (images, PDFs, etc.).
            For text files, use write_file() instead.
        """
        try:
            # Validate content type
            if not isinstance(content, (bytes, bytearray)):
                return (
                    f"Error: content must be bytes or bytearray, "
                    f"got {type(content).__name__}"
                )

            # Check size
            content_size = len(content)
            if content_size > max_size:
                return (
                    f"Error: Content too large ({content_size} bytes). "
                    f"Maximum allowed: {max_size} bytes."
                )

            resolved_path = self._resolve_path(path)

            # Create directory if needed
            dir_path = os.path.dirname(resolved_path)
            if dir_path:
                os.makedirs(dir_path, exist_ok=True)

            # Write binary content
            with open(resolved_path, 'wb') as f:
                f.write(content)

            return (
                f"Success: Binary content written to {path} "
                f"({content_size} bytes)"
            )

        except PermissionError:
            return f"Error: Permission denied when writing to {path}"
        except IsADirectoryError:
            return f"Error: {path} is a directory, not a file"
        except OSError as e:
            return f"Error: OS error when writing to {path}: {str(e)}"
        except Exception as e:
            return f"Error writing to {path}: {str(e)}"
