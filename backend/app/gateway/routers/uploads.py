"""Upload router for handling file uploads."""

import logging
from pathlib import Path
import hashlib
import json
from typing import Any

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel

from deerflow.config.paths import VIRTUAL_PATH_PREFIX, get_paths
from deerflow.sandbox.sandbox_provider import get_sandbox_provider
from deerflow.utils.file_conversion import CONVERTIBLE_EXTENSIONS, convert_file_to_markdown

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/threads/{thread_id}/uploads", tags=["uploads"])


class UploadResponse(BaseModel):
    """Response model for file upload."""

    success: bool
    files: list[dict[str, str]]
    message: str


def get_uploads_dir(thread_id: str) -> Path:
    """Get the uploads directory for a thread.

    Args:
        thread_id: The thread ID.

    Returns:
        Path to the uploads directory.
    """
    base_dir = get_paths().sandbox_uploads_dir(thread_id)
    base_dir.mkdir(parents=True, exist_ok=True)
    return base_dir

def get_file_hash(file_path: Path) -> str:
    """Calculate MD5 hash of a file."""
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()


def detect_images_in_document(file_path: Path) -> dict:
    """
    Detect if a document contains images.

    Returns:
        {
            'has_images': bool,
            'image_count': int,
            'pages_with_images': list[int],
            'metadata': dict
        }
    """
    ext = file_path.suffix.lower()

    if ext == '.pdf':
        return detect_images_in_pdf(file_path)
    elif ext in ['.docx', '.pptx']:
        return detect_images_in_office(file_path, ext)
    else:
        return {'has_images': False, 'image_count': 0}


def detect_images_in_pdf(pdf_path: Path) -> dict:
    """Detect images in a PDF file."""
    try:
        import PyPDF2

        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            pages_with_images = []
            total_images = 0

            for page_num, page in enumerate(reader.pages):
                page_images = 0

                # Check for /Resources and /XObject
                if '/Resources' in page and '/XObject' in page['/Resources']:
                    try:
                        xObject = page['/Resources']['/XObject'].get_object()
                        for obj in xObject:
                            if xObject[obj]['/Subtype'] == '/Image':
                                page_images += 1
                    except:
                        pass

                if page_images > 0:
                    pages_with_images.append(page_num + 1)
                    total_images += page_images

        return {
            'has_images': total_images > 0,
            'image_count': total_images,
            'pages_with_images': pages_with_images,
            'metadata': {
                'total_pages': len(reader.pages),
                'has_text_layer': True
            }
        }
    except Exception as e:
        logger.warning(f"Failed to detect images in PDF {pdf_path.name}: {e}")
        return {'has_images': False, 'image_count': 0}


def detect_images_in_office(file_path: Path, ext: str) -> dict:
    """Detect images in Office documents."""
    try:
        if ext == '.docx':
            from docx import Document
            doc = Document(file_path)

            image_count = 0
            # Count images in document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1

            return {
                'has_images': image_count > 0,
                'image_count': image_count,
                'metadata': {'file_type': 'docx'}
            }

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)

            image_count = 0
            pages_with_images = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_images = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        slide_images += 1

                if slide_images > 0:
                    pages_with_images.append(slide_num)
                    image_count += slide_images

            return {
                'has_images': image_count > 0,
                'image_count': image_count,
                'pages_with_images': pages_with_images,
                'metadata': {
                    'total_slides': len(prs.slides),
                    'file_type': 'pptx'
                }
            }

    except Exception as e:
        logger.warning(f"Failed to detect images in {ext} {file_path.name}: {e}")
        return {'has_images': False, 'image_count': 0}


async def convert_file_to_markdown(file_path: Path) -> Path | None:
    """
    Convert file to markdown with smart image handling.

    Strategy:
    - No images: Direct conversion with markitdown
    - Has images: Create placeholder, defer to LLM analysis on read
    """
    # Check file size
    file_size = file_path.stat().st_size
    max_size = 10 * 1024 * 1024  # 10MB
    if file_size > max_size:
        logger.warning(f"File too large: {file_path.name} ({file_size} bytes)")
        return None

    # Detect images
    image_info = detect_images_in_document(file_path)

    ext = file_path.suffix.lower()

    # For PDF with images or scanned
    if ext == '.pdf':
        # Check if scanned
        is_scanned = await detect_scanned_pdf(file_path)

        if is_scanned:
            # Create placeholder for scanned PDF
            logger.info(f"Scanned PDF detected: {file_path.name}")
            return create_scanned_pdf_placeholder(file_path)

        elif image_info['has_images']:
            # PDF with images but has text layer
            logger.info(
                f"PDF with images: {file_path.name} "
                f"({image_info['image_count']} images on pages {image_info['pages_with_images']})"
            )
            return create_pdf_with_images_placeholder(file_path, image_info)

    # For Office docs with images
    elif ext in ['.docx', '.pptx'] and image_info['has_images']:
        logger.info(
            f"Office doc with images: {file_path.name} "
            f"({image_info['image_count']} images)"
        )
        return create_office_with_images_placeholder(file_path, image_info)

    # For docs without images, use markitdown
    try:
        from markitdown import MarkItDown
        md = MarkItDown(enable_plugins=False)
        result = md.convert(str(file_path))

        if not result.text_content or len(result.text_content.strip()) < 10:
            logger.warning(f"Empty/minimal content from {file_path.name}")
            return None

        # Sanitize
        sanitized = sanitize_markdown_content(result.text_content)
        if isinstance(file_path, str):
            file_path = Path(file_path)

        md_path = file_path.with_suffix(".md")
        md_path.write_text(sanitized, encoding="utf-8")

        logger.info(f"Converted {file_path.name} to markdown (no images)")
        return md_path

    except Exception as e:
        logger.error(f"Failed to convert {file_path.name}: {e}")
        return None


async def detect_scanned_pdf(pdf_path: Path) -> bool:
    """Detect if a PDF is a scanned document."""
    try:
        import PyPDF2

        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)

            # Check first 3 pages
            for page in reader.pages[:3]:
                text = page.extract_text()
                if text and len(text.strip()) > 50:
                    # Found meaningful text
                    return False

            return True

    except Exception as e:
        logger.warning(f"Failed to detect scanned PDF: {e}")
        return False


def create_scanned_pdf_placeholder(pdf_path: Path) -> Path:
    """Create placeholder for scanned PDF."""
    placeholder = f"""# 📄 Scanned PDF Document

**File**: `{pdf_path.name}`
**Size**: {pdf_path.stat().st_size} bytes
**Type**: Scanned (no extractable text)

> This document appears to be a scanned PDF. The content will be extracted 
> using AI OCR when you read this file.

## How to Read This Document

Use the `read_file` tool - it will automatically:
1. Convert each page to an image
2. Use an AI vision model to extract and understand the text
3. Provide detailed analysis including:
   - Extracted text (with high accuracy)
   - Image descriptions and understanding
   - Table recognition and data extraction
   - Chart and graph interpretation

**Note**: This will use a vision-capable AI model (GPT-4V, Claude-3.5 Sonnet, etc.)
"""
    if isinstance(pdf_path, str):
        pdf_path = Path(pdf_path)
    md_path = pdf_path.with_suffix(".md")
    md_path.write_text(placeholder, encoding="utf-8")

    # Save metadata for later use
    save_file_metadata(pdf_path, {'type': 'scanned_pdf', 'needs_ocr': True})

    return md_path


def create_pdf_with_images_placeholder(pdf_path: Path, image_info: dict) -> Path:
    """Create placeholder for PDF with images."""
    placeholder = f"""# 📄 PDF Document with Images

**File**: `{pdf_path.name}`
**Total Pages**: {image_info.get('metadata', {}).get('total_pages', 'N/A')}
**Pages with Images**: {', '.join(map(str, image_info.get('pages_with_images', [])))}
**Total Images**: {image_info['image_count']}

> This document contains both text and images. When you read this file, 
> the AI will provide:
> 
> 1. Text content (extracted directly)
> 2. Detailed analysis of each embedded image
> 3. Understanding of charts, graphs, diagrams
> 4. Table recognition and data extraction

## Content Preview

The full content including image analysis will be loaded when you use the `read_file` tool.

---

### Images in This Document

"""

    pages_with_images = image_info.get('pages_with_images', [])
    for page_num in pages_with_images:
        placeholder += f"- **Page {page_num}**: Contains images\n"

    placeholder += "\n\nUse `read_file` to get the complete content with image analysis."
    if isinstance(pdf_path, str):
        pdf_path = Path(pdf_path)
    md_path = pdf_path.with_suffix(".md")
    md_path.write_text(placeholder, encoding="utf-8")

    # Save metadata
    save_file_metadata(pdf_path, {
        'type': 'pdf_with_images',
        'needs_image_analysis': True,
        'pages_with_images': pages_with_images,
        'image_count': image_info['image_count']
    })

    return md_path


def create_office_with_images_placeholder(file_path: Path, image_info: dict) -> Path:
    """Create placeholder for Office documents with images."""
    file_type = image_info.get('metadata', {}).get('file_type', 'document')

    placeholder = f"""# 📄 {file_type.upper()} Document with Images

**File**: `{file_path.name}`
**Total Images**: {image_info['image_count']}
"""

    if file_type == 'pptx':
        total_slides = image_info.get('metadata', {}).get('total_slides', 'N/A')
        slides_with_images = image_info.get('pages_with_images', [])

        placeholder += f"""**Total Slides**: {total_slides}
**Slides with Images**: {', '.join(map(str, slides_with_images))}
"""

        placeholder += """
This presentation contains images. When you read this file, the AI will provide:
1. Slide text content
2. Detailed analysis of each embedded image
3. Understanding of charts, diagrams, screenshots
4. Contextual interpretation of visual elements

"""

        if slides_with_images:
            placeholder += "### Slides with Images\n\n"
            for slide_num in slides_with_images:
                placeholder += f"- **Slide {slide_num}**: Contains images\n"

    elif file_type == 'docx':
        placeholder += """
This document contains images. When you read this file, the AI will provide:
1. Document text content
2. Detailed analysis of each embedded image
3. Understanding of figures, screenshots, charts
4. Contextual integration of visual elements

"""

    placeholder += "\nUse `read_file` to get the complete content with image analysis."
    if isinstance(file_path, str):
        file_path = Path(file_path)

    md_path = file_path.with_suffix(".md")
    md_path.write_text(placeholder, encoding="utf-8")

    # Save metadata
    save_file_metadata(file_path, {
        'type': f'{file_type}_with_images',
        'needs_image_analysis': True,
        'image_count': image_info['image_count']
    })

    return md_path


def save_file_metadata(file_path: Path, metadata: dict) -> None:
    """Save metadata for a file."""
    if isinstance(file_path, str):
        file_path = Path(file_path)

    meta_path = file_path.with_suffix('.meta.json')
    with open(meta_path, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)


def load_file_metadata(file_path: Path) -> dict:
    """Load metadata for a file."""
    if isinstance(file_path, str):
        file_path = Path(file_path)

    meta_path = file_path.with_suffix('.meta.json')
    if meta_path.exists():
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            pass
    return {}


def sanitize_markdown_content(content: str, max_length: int = 100000) -> str:
    """Sanitize markdown content to prevent encoding errors."""
    if not content:
        return ""

    import re

    # Remove control characters except newlines and tabs
    content = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', content)

    # Fix Unicode
    content = content.encode('utf-8', errors='replace').decode('utf-8', errors='replace')

    # Normalize whitespace
    content = re.sub(r'[ \t]+', ' ', content)
    content = re.sub(r'\n{3,}', '\n\n', content)

    # Truncate
    if len(content) > max_length:
        content = content[:max_length] + (
            f"\n\n[Content truncated from {len(content)} to {max_length} characters]"
        )

    return content.strip()

# async def convert_file_to_markdown(file_path: Path) -> Path | None:
#     """Convert a file to markdown using markitdown.

#     Args:
#         file_path: Path to the file to convert.

#     Returns:
#         Path to the markdown file if conversion was successful, None otherwise.
#     """
#     try:
#         from markitdown import MarkItDown

#         md = MarkItDown()
#         result = md.convert(str(file_path))

#         # Save as .md file with same name
#         md_path = file_path.with_suffix(".md")
#         md_path.write_text(result.text_content, encoding="utf-8", errors="replace")

#         logger.info(f"Converted {file_path.name} to markdown: {md_path.name}")
#         return md_path
#     except Exception as e:
#         logger.error(f"Failed to convert {file_path.name} to markdown: {e}")
#         return None
    

@router.post("", response_model=UploadResponse)
async def upload_files(
    thread_id: str,
    files: list[UploadFile] = File(...),
) -> UploadResponse:
    """Upload multiple files to a thread's uploads directory.

    For PDF, PPT, Excel, and Word files, they will be converted to markdown using markitdown.
    All files (original and converted) are saved to /mnt/user-data/uploads.

    Args:
        thread_id: The thread ID to upload files to.
        files: List of files to upload.

    Returns:
        Upload response with success status and file information.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    uploads_dir = get_uploads_dir(thread_id)
    paths = get_paths()
    uploaded_files = []

    sandbox_provider = get_sandbox_provider()
    sandbox_id = sandbox_provider.acquire(thread_id)
    sandbox = sandbox_provider.get(sandbox_id)

    for file in files:
        if not file.filename:
            continue

        try:
            # Normalize filename to prevent path traversal
            safe_filename = Path(file.filename).name
            if not safe_filename or safe_filename in {".", ".."} or "/" in safe_filename or "\\" in safe_filename:
                logger.warning(f"Skipping file with unsafe filename: {file.filename!r}")
                continue

            content = await file.read()
            file_path = uploads_dir / safe_filename
            file_path.write_bytes(content)

            # Build relative path from backend root
            relative_path = str(paths.sandbox_uploads_dir(thread_id) / safe_filename)
            virtual_path = f"{VIRTUAL_PATH_PREFIX}/uploads/{safe_filename}"

            # Keep local sandbox source of truth in thread-scoped host storage.
            # For non-local sandboxes, also sync to virtual path for runtime visibility.
            if sandbox_id != "local":
                sandbox.update_file(virtual_path, content)

            file_info = {
                "filename": safe_filename,
                "size": str(len(content)),
                "path": relative_path,  # Actual filesystem path (relative to backend/)
                "virtual_path": virtual_path,  # Path for Agent in sandbox
                "artifact_url": f"/api/threads/{thread_id}/artifacts/mnt/user-data/uploads/{safe_filename}",  # HTTP URL
            }

            logger.info(f"Saved file: {safe_filename} ({len(content)} bytes) to {relative_path}")

            # Check if file should be converted to markdown
            file_ext = file_path.suffix.lower()
            if file_ext in CONVERTIBLE_EXTENSIONS:
                md_path = await convert_file_to_markdown(file_path)
                if md_path:
                    md_relative_path = str(paths.sandbox_uploads_dir(thread_id) / md_path.name)
                    md_virtual_path = f"{VIRTUAL_PATH_PREFIX}/uploads/{md_path.name}"

                    if sandbox_id != "local":
                        sandbox.update_file(md_virtual_path, md_path.read_bytes())

                    file_info["markdown_file"] = md_path.name
                    file_info["markdown_path"] = md_relative_path
                    file_info["markdown_virtual_path"] = md_virtual_path
                    file_info["markdown_artifact_url"] = f"/api/threads/{thread_id}/artifacts/mnt/user-data/uploads/{md_path.name}"

            uploaded_files.append(file_info)

        except Exception as e:
            logger.error(f"Failed to upload {file.filename}: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to upload {file.filename}: {str(e)}")

    return UploadResponse(
        success=True,
        files=uploaded_files,
        message=f"Successfully uploaded {len(uploaded_files)} file(s)",
    )


@router.get("/list", response_model=dict)
async def list_uploaded_files(thread_id: str) -> dict:
    """List all files in a thread's uploads directory.

    Args:
        thread_id: The thread ID to list files for.

    Returns:
        Dictionary containing list of files with their metadata.
    """
    uploads_dir = get_uploads_dir(thread_id)

    if not uploads_dir.exists():
        return {"files": [], "count": 0}

    files = []
    for file_path in sorted(uploads_dir.iterdir()):
        if file_path.is_file():
            stat = file_path.stat()
            relative_path = str(get_paths().sandbox_uploads_dir(thread_id) / file_path.name)
            files.append(
                {
                    "filename": file_path.name,
                    "size": stat.st_size,
                    "path": relative_path,  # Actual filesystem path
                    "virtual_path": f"{VIRTUAL_PATH_PREFIX}/uploads/{file_path.name}",  # Path for Agent in sandbox
                    "artifact_url": f"/api/threads/{thread_id}/artifacts/mnt/user-data/uploads/{file_path.name}",  # HTTP URL
                    "extension": file_path.suffix,
                    "modified": stat.st_mtime,
                }
            )

    return {"files": files, "count": len(files)}


@router.delete("/{filename}")
async def delete_uploaded_file(thread_id: str, filename: str) -> dict:
    """Delete a file from a thread's uploads directory.

    Args:
        thread_id: The thread ID.
        filename: The filename to delete.

    Returns:
        Success message.
    """
    uploads_dir = get_uploads_dir(thread_id)
    file_path = uploads_dir / filename

    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"File not found: {filename}")

    # Security check: ensure the path is within the uploads directory
    try:
        file_path.resolve().relative_to(uploads_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")

    try:
        file_path.unlink()
        logger.info(f"Deleted file: {filename}")
        return {"success": True, "message": f"Deleted {filename}"}
    except Exception as e:
        logger.error(f"Failed to delete {filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete {filename}: {str(e)}")
