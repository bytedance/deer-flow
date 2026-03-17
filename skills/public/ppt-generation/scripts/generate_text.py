"""Generate a text-only PowerPoint presentation from a plan JSON file."""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_COLOR = RGBColor(0x00, 0xB4, 0xD8)
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
BODY_COLOR = RGBColor(0xCC, 0xCC, 0xCC)


def set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TITLE_COLOR

    slide.shapes.add_connector(1, Inches(1), Inches(4.2), Inches(5), Inches(4.2))

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.33), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle
        subtitle_paragraph.font.size = Pt(24)
        subtitle_paragraph.font.color.rgb = ACCENT_COLOR


def add_content_slide(prs: Presentation, title: str, key_points: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.9))
    title_frame = title_box.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = ACCENT_COLOR

    slide.shapes.add_connector(1, Inches(0.5), Inches(1.35), Inches(12.83), Inches(1.35))

    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    for index, point in enumerate(key_points):
        paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        paragraph.text = f"- {point}"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = BODY_COLOR
        paragraph.space_after = Pt(8)


def generate_ppt(plan_file: str, output_file: str) -> str:
    with open(plan_file, "r", encoding="utf-8-sig") as file:
        plan = json.load(file)

    os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_info in plan.get("slides", []):
        slide_type = slide_info.get("type", "content")
        title = slide_info.get("title", "")
        subtitle = slide_info.get("subtitle", "")
        key_points = slide_info.get("key_points", [])

        if slide_type == "title" or (not key_points and subtitle):
            add_title_slide(prs, title, subtitle)
        else:
            add_content_slide(prs, title, key_points)

    prs.save(output_file)
    return f"Successfully generated presentation with {len(plan.get('slides', []))} slides to {output_file}"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate text-based PPTX from plan JSON")
    parser.add_argument("--plan-file", required=True, help="Path to JSON presentation plan file")
    parser.add_argument("--output-file", required=True, help="Output path for PPTX file")
    args = parser.parse_args()

    try:
        print(generate_ppt(args.plan_file, args.output_file))
    except Exception as error:
        print(f"Error: {error}", file=sys.stderr)
        sys.exit(1)
