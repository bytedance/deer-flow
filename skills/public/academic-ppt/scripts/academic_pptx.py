#!/usr/bin/env python3
"""Academic PPTX generator for DeerFlow.

Generates native editable PowerPoint presentations with academic styling,
LaTeX formula rendering, chart embedding, and speaker notes.
"""

import argparse
import json
import os
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


STYLES = {
    "clean-academic": {
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x1A, 0x2E),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x2E, 0x86, 0xC1),
        "header_bg": None,
        "title_font": "Georgia",
        "body_font": "Calibri",
        "title_size": Pt(32),
        "body_size": Pt(18),
        "header_height": Inches(0),
    },
    "beamer-blue": {
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x2C, 0x3E, 0x50),
        "accent_color": RGBColor(0x2C, 0x3E, 0x80),
        "header_bg": RGBColor(0x2C, 0x3E, 0x80),
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": Pt(28),
        "body_size": Pt(18),
        "header_height": Inches(1.2),
    },
    "beamer-red": {
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x2C, 0x2C, 0x2C),
        "accent_color": RGBColor(0x8B, 0x00, 0x00),
        "header_bg": RGBColor(0x8B, 0x00, 0x00),
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": Pt(28),
        "body_size": Pt(18),
        "header_height": Inches(1.2),
    },
    "thesis-defense": {
        "bg_color": RGBColor(0xFA, 0xFA, 0xFA),
        "title_color": RGBColor(0x1B, 0x2A, 0x4A),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "accent_color": RGBColor(0x1B, 0x2A, 0x4A),
        "header_bg": None,
        "title_font": "Times New Roman",
        "body_font": "Calibri",
        "title_size": Pt(30),
        "body_size": Pt(18),
        "header_height": Inches(0),
    },
    "research-meeting": {
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x2D, 0x2D, 0x2D),
        "text_color": RGBColor(0x44, 0x44, 0x44),
        "accent_color": RGBColor(0x00, 0x7A, 0xCC),
        "header_bg": None,
        "title_font": "Arial",
        "body_font": "Arial",
        "title_size": Pt(26),
        "body_size": Pt(16),
        "header_height": Inches(0),
    },
    "poster-style": {
        "bg_color": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x00, 0x3D, 0x6B),
        "text_color": RGBColor(0x1A, 0x1A, 0x1A),
        "accent_color": RGBColor(0x00, 0x3D, 0x6B),
        "header_bg": None,
        "title_font": "Arial Black",
        "body_font": "Arial",
        "title_size": Pt(36),
        "body_size": Pt(20),
        "header_height": Inches(0),
    },
}


def create_presentation(plan: dict, style_name: str = "clean-academic", widescreen: bool = True) -> Presentation:
    """Create a PPTX presentation from a structured plan."""
    style = STYLES.get(style_name, STYLES["clean-academic"])
    prs = Presentation()

    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    total_slides = len(plan.get("slides", []))

    for slide_data in plan.get("slides", []):
        slide_type = slide_data.get("type", "content")
        slide_num = slide_data.get("slide_number", 0)

        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        _set_slide_background(slide, style["bg_color"])

        if slide_type == "title":
            _create_title_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "outline":
            _create_outline_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "section":
            _create_section_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "content":
            _create_content_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "figure":
            _create_figure_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "formula":
            _create_formula_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "table":
            _create_table_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "two_column":
            _create_two_column_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "references":
            _create_references_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "thank_you":
            _create_thank_you_slide(slide, slide_data, style, slide_width, slide_height)
        elif slide_type == "assertion_evidence":
            _create_assertion_evidence_slide(slide, slide_data, style, slide_width, slide_height)
        else:
            _create_content_slide(slide, slide_data, style, slide_width, slide_height)

        if slide_type not in ("title", "thank_you") and slide_num > 0:
            _add_slide_number(slide, slide_num, total_slides, style, slide_width, slide_height)

        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    return prs


def _set_slide_background(slide, color: RGBColor):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, style: dict, slide_width, title_text: str):
    """Add a colored header bar (Beamer style)."""
    if style["header_bg"] is None:
        return Inches(0.5)

    header_height = style["header_height"]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, header_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = style["header_bg"]
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), slide_width - Inches(1), header_height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = style["title_size"]
    p.font.color.rgb = style["title_color"]
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    return header_height + Inches(0.3)


def _add_title_text(slide, title: str, style: dict, slide_width, top: int = None):
    """Add a title text with underline accent."""
    if top is None:
        top = Inches(0.4)

    txBox = slide.shapes.add_textbox(Inches(0.5), top, slide_width - Inches(1), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = style["title_size"]
    p.font.color.rgb = style["title_color"] if style["header_bg"] is None else style["accent_color"]
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    line_top = top + Inches(0.7)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), line_top, Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = style["accent_color"]
    line.line.fill.background()

    return line_top + Inches(0.3)


def _create_assertion_evidence_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create an assertion-evidence slide: full-sentence title + visual evidence (no bullets)."""
    assertion = data.get("assertion", data.get("title", ""))
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = assertion
    p.font.size = Pt(24)
    p.font.color.rgb = style.get("accent_color", style["text_color"])
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    figure_path = data.get("figure_path", "")
    if figure_path and os.path.exists(figure_path):
        max_w = slide_width - Inches(1.2)
        max_h = slide_height - Inches(2.8)
        try:
            from PIL import Image
            img = Image.open(figure_path)
            w_ratio = max_w / Emu(int(img.width * 914400 / 96))
            h_ratio = max_h / Emu(int(img.height * 914400 / 96))
            ratio = min(w_ratio, h_ratio, 1.0)
            pic_w = int(img.width * 914400 / 96 * ratio)
            pic_h = int(img.height * 914400 / 96 * ratio)
        except ImportError:
            pic_w, pic_h = max_w, max_h
        left = (slide_width - pic_w) // 2
        slide.shapes.add_picture(figure_path, left, Inches(1.6), pic_w, pic_h)

    caption = data.get("caption", "")
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.6), slide_height - Inches(0.8), slide_width - Inches(1.2), Inches(0.5))
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_p = cap_tf.paragraphs[0]
        cap_p.text = caption
        cap_p.font.size = Pt(12)
        cap_p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        cap_p.font.name = style["body_font"]
        cap_p.font.italic = True
        cap_p.alignment = PP_ALIGN.CENTER


def _create_title_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a title slide."""
    center_y = slide_height // 2

    if style["header_bg"]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = style["header_bg"]
        shape.line.fill.background()
        title_color = style["title_color"]
        text_color = RGBColor(0xDD, 0xDD, 0xDD)
    else:
        title_color = style["accent_color"]
        text_color = style["text_color"]

    txBox = slide.shapes.add_textbox(Inches(1), center_y - Inches(1.5), slide_width - Inches(2), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(36)
    p.font.color.rgb = title_color
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle = data.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = text_color
        p2.font.name = style["body_font"]
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    authors = data.get("authors", "")
    if authors:
        txBox2 = slide.shapes.add_textbox(Inches(1), center_y + Inches(0.5), slide_width - Inches(2), Inches(0.5))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = authors
        p3.font.size = Pt(16)
        p3.font.color.rgb = text_color
        p3.font.name = style["body_font"]
        p3.alignment = PP_ALIGN.CENTER

    affiliations = data.get("affiliations", "")
    if affiliations:
        txBox3 = slide.shapes.add_textbox(Inches(1), center_y + Inches(1.0), slide_width - Inches(2), Inches(0.4))
        tf3 = txBox3.text_frame
        p4 = tf3.paragraphs[0]
        p4.text = affiliations
        p4.font.size = Pt(13)
        p4.font.color.rgb = text_color
        p4.font.name = style["body_font"]
        p4.font.italic = True
        p4.alignment = PP_ALIGN.CENTER

    date_text = data.get("date", "")
    if date_text:
        txBox4 = slide.shapes.add_textbox(Inches(1), center_y + Inches(1.5), slide_width - Inches(2), Inches(0.4))
        tf4 = txBox4.text_frame
        p5 = tf4.paragraphs[0]
        p5.text = date_text
        p5.font.size = Pt(14)
        p5.font.color.rgb = text_color
        p5.font.name = style["body_font"]
        p5.alignment = PP_ALIGN.CENTER


def _create_outline_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create an outline/TOC slide."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", "Outline"))
    else:
        content_top = _add_title_text(slide, data.get("title", "Outline"), style, slide_width)

    items = data.get("items", [])
    txBox = slide.shapes.add_textbox(Inches(1.5), content_top + Inches(0.3), slide_width - Inches(3), slide_height - content_top - Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {i + 1}.  {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = style["text_color"]
        p.font.name = style["body_font"]
        p.space_before = Pt(16)
        p.space_after = Pt(8)


def _create_section_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a section divider slide."""
    if style["header_bg"]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = style["header_bg"]
        shape.line.fill.background()
        color = style["title_color"]
    else:
        color = style["accent_color"]

    center_y = slide_height // 2
    txBox = slide.shapes.add_textbox(Inches(1), center_y - Inches(0.5), slide_width - Inches(2), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(36)
    p.font.color.rgb = color
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def _create_content_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a standard content slide with bullets."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", ""))
    else:
        content_top = _add_title_text(slide, data.get("title", ""), style, slide_width)

    bullets = data.get("bullets", [])
    txBox = slide.shapes.add_textbox(Inches(0.8), content_top + Inches(0.2), slide_width - Inches(1.6), slide_height - content_top - Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = style["body_size"]
        p.font.color.rgb = style["text_color"]
        p.font.name = style["body_font"]
        p.space_before = Pt(10)
        p.space_after = Pt(6)


def _create_figure_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a slide with an embedded figure."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", ""))
    else:
        content_top = _add_title_text(slide, data.get("title", ""), style, slide_width)

    figure_path = data.get("figure_path", "")
    if figure_path and os.path.exists(figure_path):
        available_height = slide_height - content_top - Inches(1.2)
        available_width = slide_width - Inches(2)

        try:
            from PIL import Image
            img = Image.open(figure_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            if available_width / aspect <= available_height:
                w = available_width
                h = int(w / aspect)
            else:
                h = available_height
                w = int(h * aspect)

            left = (slide_width - w) // 2
            slide.shapes.add_picture(figure_path, left, content_top + Inches(0.2), w, h)
        except ImportError:
            slide.shapes.add_picture(figure_path, Inches(1), content_top + Inches(0.2), slide_width - Inches(2), available_height)
    else:
        txBox = slide.shapes.add_textbox(Inches(1), content_top + Inches(1), slide_width - Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Figure: {figure_path}]"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER

    caption = data.get("caption", "")
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(1), slide_height - Inches(0.8), slide_width - Inches(2), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER


def _create_formula_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a slide with a LaTeX formula."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", ""))
    else:
        content_top = _add_title_text(slide, data.get("title", ""), style, slide_width)

    formula_latex = data.get("formula_latex", "")
    formula_image = _render_latex_to_image(formula_latex)

    if formula_image and os.path.exists(formula_image):
        try:
            from PIL import Image
            img = Image.open(formula_image)
            img_w, img_h = img.size
            max_width = slide_width - Inches(3)
            scale = min(1.0, max_width / Emu(img_w * 914400 / 300))
            w = Inches(img_w / 300 * scale)
            h = Inches(img_h / 300 * scale)
            left = (slide_width - w) // 2
            top = content_top + Inches(1)
            slide.shapes.add_picture(formula_image, left, top, w, h)
        except Exception:
            left = Inches(1)
            top = content_top + Inches(1)
            slide.shapes.add_picture(formula_image, left, top)
    else:
        txBox = slide.shapes.add_textbox(Inches(1), content_top + Inches(1), slide_width - Inches(2), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"${formula_latex}$"
        p.font.size = Pt(24)
        p.font.color.rgb = style["text_color"]
        p.font.name = "Cambria Math"
        p.alignment = PP_ALIGN.CENTER

    explanation = data.get("explanation", "")
    if explanation:
        exp_box = slide.shapes.add_textbox(Inches(1), slide_height - Inches(2), slide_width - Inches(2), Inches(0.8))
        tf = exp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = explanation
        p.font.size = Pt(16)
        p.font.color.rgb = style["text_color"]
        p.font.name = style["body_font"]
        p.alignment = PP_ALIGN.CENTER


def _create_table_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a slide with a data table."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", ""))
    else:
        content_top = _add_title_text(slide, data.get("title", ""), style, slide_width)

    headers = data.get("headers", [])
    rows = data.get("rows", [])

    if not headers or not rows:
        return

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_width = min(slide_width - Inches(1.5), Inches(n_cols * 2))
    table_height = Inches(n_rows * 0.45)
    left = (slide_width - table_width) // 2

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, content_top + Inches(0.4), table_width, table_height)
    table = table_shape.table

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.font.name = style["body_font"]
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = style["accent_color"]

    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            clean_text = cell_text.replace("**", "")
            cell.text = clean_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = style["text_color"]
                paragraph.font.name = style["body_font"]
                paragraph.alignment = PP_ALIGN.CENTER
                if "**" in cell_text:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = style["accent_color"]

            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)


def _create_two_column_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a two-column comparison slide."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", ""))
    else:
        content_top = _add_title_text(slide, data.get("title", ""), style, slide_width)

    col_width = (slide_width - Inches(1.5)) // 2

    for col_idx, side in enumerate(["left", "right"]):
        left = Inches(0.5) + col_idx * (col_width + Inches(0.5))
        col_title = data.get(f"{side}_title", "")

        title_box = slide.shapes.add_textbox(left, content_top + Inches(0.2), col_width, Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = style["accent_color"]
        p.font.name = style["body_font"]

        bullets = data.get(f"{side}_bullets", [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(left, content_top + Inches(0.9), col_width, slide_height - content_top - Inches(1.5))
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = style["text_color"]
                p.font.name = style["body_font"]
                p.space_before = Pt(8)


def _create_references_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a references/bibliography slide."""
    if style["header_bg"]:
        content_top = _add_header_bar(slide, style, slide_width, data.get("title", "References"))
    else:
        content_top = _add_title_text(slide, data.get("title", "References"), style, slide_width)

    references = data.get("references", [])
    txBox = slide.shapes.add_textbox(Inches(0.5), content_top + Inches(0.2), slide_width - Inches(1), slide_height - content_top - Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(references):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(12)
        p.font.color.rgb = style["text_color"]
        p.font.name = style["body_font"]
        p.space_before = Pt(4)
        p.space_after = Pt(4)


def _create_thank_you_slide(slide, data: dict, style: dict, slide_width, slide_height):
    """Create a thank-you / closing slide."""
    if style["header_bg"]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = style["header_bg"]
        shape.line.fill.background()
        color = style["title_color"]
        sub_color = RGBColor(0xCC, 0xCC, 0xCC)
    else:
        color = style["accent_color"]
        sub_color = style["text_color"]

    center_y = slide_height // 2
    txBox = slide.shapes.add_textbox(Inches(1), center_y - Inches(1), slide_width - Inches(2), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "Thank You!")
    p.font.size = Pt(40)
    p.font.color.rgb = color
    p.font.name = style["title_font"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle = data.get("subtitle", "Questions?")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = sub_color
        p2.font.name = style["body_font"]
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    contact = data.get("contact", "")
    if contact:
        p3 = tf.add_paragraph()
        p3.text = contact
        p3.font.size = Pt(16)
        p3.font.color.rgb = sub_color
        p3.font.name = style["body_font"]
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(30)


def _add_slide_number(slide, num: int, total: int, style: dict, slide_width, slide_height):
    """Add slide number to bottom right."""
    txBox = slide.shapes.add_textbox(slide_width - Inches(1.2), slide_height - Inches(0.4), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.font.name = style["body_font"]
    p.alignment = PP_ALIGN.RIGHT


def _render_latex_to_image(latex: str) -> str | None:
    """Render LaTeX formula to a PNG image using matplotlib."""
    if not latex:
        return None
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(8, 1.5))
        ax.axis("off")
        ax.text(0.5, 0.5, f"${latex}$", fontsize=28, ha="center", va="center",
                transform=ax.transAxes, usetex=False)
        fig.tight_layout(pad=0.1)

        tmp_path = tempfile.mktemp(suffix=".png")
        fig.savefig(tmp_path, dpi=300, bbox_inches="tight", transparent=True, pad_inches=0.1)
        plt.close(fig)
        return tmp_path
    except Exception as e:
        print(f"Warning: Could not render LaTeX formula: {e}", file=sys.stderr)
        return None


def main():
    parser = argparse.ArgumentParser(description="Academic PPTX Generator")
    parser.add_argument("--plan-file", required=True, help="Path to presentation plan JSON")
    parser.add_argument("--output-file", required=True, help="Output PPTX file path")
    parser.add_argument("--style", default="clean-academic", choices=list(STYLES.keys()), help="Presentation style")
    parser.add_argument("--widescreen", default=True, type=lambda x: x.lower() != "false", help="16:9 widescreen")
    args = parser.parse_args()

    with open(args.plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    style = args.style or plan.get("style", "clean-academic")
    prs = create_presentation(plan, style_name=style, widescreen=args.widescreen)

    os.makedirs(os.path.dirname(os.path.abspath(args.output_file)), exist_ok=True)
    prs.save(args.output_file)

    n_slides = len(plan.get("slides", []))
    print(f"Academic presentation generated successfully!")
    print(f"  Style: {style}")
    print(f"  Slides: {n_slides}")
    print(f"  Output: {args.output_file}")


if __name__ == "__main__":
    main()
